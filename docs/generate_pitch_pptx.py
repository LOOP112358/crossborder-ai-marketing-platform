# -*- coding: utf-8 -*-
"""生成「越境智绘」项目展示汇报 PPT（蓝色科技 · 简洁专业）"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
from pathlib import Path

OUT = Path(__file__).resolve().parent / "YueJingZhiHui-project-pitch-tech.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Professional tech blue
INK = RGBColor(0x1A, 0x23, 0x3A)
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
BLUE = RGBColor(0x1E, 0x5A, 0xA8)
BLUE_MID = RGBColor(0x2F, 0x6F, 0xDB)
CYAN = RGBColor(0x1A, 0xB4, 0xD8)
BG = RGBColor(0xF4, 0xF7, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x5B, 0x6B, 0x7C)
LINE = RGBColor(0xD7, 0xE0, 0xEC)
CARD = WHITE
FONT = "Microsoft YaHei"


def set_run(run, size=18, bold=False, color=INK, font=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    r_pr = run._r.get_or_add_rPr()
    for child in list(r_pr):
        if child.tag.endswith("}ea") or child.tag.endswith("}cs"):
            r_pr.remove(child)
    ea = etree.SubElement(r_pr, "{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    ea.set("typeface", font)


def add_rect(slide, l, t, w, h, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def add_round(slide, l, t, w, h, fill, radius=0.06):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    try:
        shape.adjustments[0] = radius
    except Exception:
        pass
    return shape


def add_text(slide, l, t, w, h, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        set_run(run, size=size, bold=bold, color=color)
    return box


def add_bullets(slide, l, t, w, h, items, size=15, color=INK):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = "•  " + item
        set_run(run, size=size, color=color)
    return box


def slide_chrome(slide, title, subtitle=None):
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, BG)
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.08), CYAN)
    add_rect(slide, 0, Inches(6.95), prs.slide_width, Inches(0.55), NAVY)
    add_text(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.35),
             "越境智绘 · 跨境AI营销系统", size=12, color=WHITE)
    add_text(slide, Inches(9.2), Inches(7.05), Inches(3.8), Inches(0.35),
             "技术答辩 · 项目展示", size=12, color=WHITE, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.5), Inches(0.35), Inches(12), Inches(0.5),
             title, size=26, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.9), Inches(12), Inches(0.35),
                 subtitle, size=13, color=MUTED)
    add_rect(slide, Inches(0.5), Inches(1.3), Inches(1.6), Inches(0.05), BLUE_MID)


def member_card(slide, x, y, w, h, title, body):
    add_round(slide, x, y, w, h, WHITE)
    add_rect(slide, x, y, Inches(0.12), h, BLUE_MID)
    add_text(slide, x + Inches(0.35), y + Inches(0.2), w - Inches(0.55), Inches(0.4),
             title, size=15, bold=True, color=BLUE)
    add_text(slide, x + Inches(0.35), y + Inches(0.7), w - Inches(0.55), h - Inches(0.9),
             body, size=12, color=INK)


# ===== 1 Cover =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, NAVY)
add_rect(s, 0, 0, Inches(0.18), prs.slide_height, CYAN)
add_rect(s, Inches(0.18), Inches(5.85), prs.slide_width - Inches(0.18), Inches(1.65), BLUE)
add_text(s, Inches(0.8), Inches(1.8), Inches(11), Inches(0.4),
         "TECHNICAL DEFENSE  ·  PROJECT PITCH", size=14, color=CYAN)
add_text(s, Inches(0.8), Inches(2.4), Inches(11.5), Inches(0.9),
         "跨境AI营销系统", size=42, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(3.4), Inches(11), Inches(0.4),
         "越境智绘 · Beyond Borders", size=18, color=RGBColor(0xA8, 0xC4, 0xE8))
add_text(s, Inches(0.8), Inches(4.1), Inches(11), Inches(0.8),
         "选品 → 文案 → 抠图 / 背景 / 海报 → 客服与看板\n一站式 AI 电商营销工具平台",
         size=16, color=RGBColor(0xC5, 0xD5, 0xE8))
add_text(s, Inches(0.9), Inches(6.15), Inches(8), Inches(0.35),
         "项目展示汇报", size=18, bold=True, color=WHITE)
add_text(s, Inches(0.9), Inches(6.6), Inches(10), Inches(0.35),
         "模块化架构 · 全链路打通 · 可演示可部署", size=13, color=RGBColor(0xC8, 0xDC, 0xF0))

# ===== 2 Why =====
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_chrome(s, "为什么做 · 服务谁", "跨境卖家素材链路长，工具分散，需要端到端闭环")
cards = [
    ("痛点", "选品、抠图、场景、海报文案、\n客服、数据分散\n素材制作链路长、协作成本高"),
    ("用户", "跨境电商商家 / 运营\n设计与内容协作方"),
    ("价值", "从 ABO 货盘到可发布营销海报\n一站式完成营销素材与运营辅助"),
]
for i, (h, body) in enumerate(cards):
    x = Inches(0.5 + i * 4.15)
    add_round(s, x, Inches(1.7), Inches(3.9), Inches(4.5), WHITE)
    add_rect(s, x, Inches(1.7), Inches(3.9), Inches(0.65), BLUE if i != 1 else NAVY)
    add_text(s, x + Inches(0.25), Inches(1.82), Inches(3.4), Inches(0.4), h, size=18, bold=True, color=WHITE)
    add_text(s, x + Inches(0.25), Inches(2.7), Inches(3.4), Inches(3), body, size=15, color=INK)

# ===== 3 Architecture =====
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_chrome(s, "整体架构一览", "统一入口 · 模块化挂载 · 本地知识库与静态资产")
blocks = [
    ("前端", "Vue 3 + Vite + Element Plus\nPinia / Router / i18n\n@m1–@m5 聚合各模块页面"),
    ("后端", "FastAPI + Uvicorn\npython run.py 统一启动\n自动挂载 module1–5 路由"),
    ("数据与资产", "SQLite · JWT 鉴权\nstatic/{matte,background,poster}\nABO 货盘 + FAISS 知识库"),
]
for i, (h, body) in enumerate(blocks):
    y = Inches(1.65 + i * 1.55)
    add_round(s, Inches(0.5), y, Inches(12.3), Inches(1.4), WHITE)
    add_rect(s, Inches(0.5), y, Inches(2.2), Inches(1.4), BLUE if i % 2 == 0 else NAVY)
    add_text(s, Inches(0.7), y + Inches(0.45), Inches(1.8), Inches(0.5), h, size=18, bold=True, color=WHITE)
    add_text(s, Inches(3.0), y + Inches(0.3), Inches(9.4), Inches(1.0), body, size=15, color=INK)

# ===== 4 Modules =====
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_chrome(s, "五大模块职责", "成员协作边界清晰，能力可独立使用也可串联")
mods = [
    ("M1 成员1", "我的工作\n（详见分工页）"),
    ("M2 抠图", "上传/识别/自动抠图\n标准化透明商品素材"),
    ("M3 背景", "Prompt 转换\nSeedream 场景背景"),
    ("M4 海报", "模板+叠字合成\n历史/收藏/下载"),
    ("M5 客服·看板·选品", "RAG 客服 / 运营看板\n选品衔接 · 云端部署"),
]
for i, (h, body) in enumerate(mods):
    x = Inches(0.35 + i * 2.55)
    add_round(s, x, Inches(1.65), Inches(2.4), Inches(4.7), WHITE)
    add_rect(s, x, Inches(1.65), Inches(2.4), Inches(0.95), BLUE if i % 2 == 0 else NAVY)
    add_text(s, x + Inches(0.1), Inches(1.85), Inches(2.2), Inches(0.6), h, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), Inches(2.9), Inches(2.1), Inches(3), body, size=13, color=INK, align=PP_ALIGN.CENTER)

# ===== 5 Workflow =====
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_chrome(s, "核心链路：AI 海报工作流", "四步推进 · 中间结果自动带入 · 底图可复用")
steps = [
    ("01", "商品抠图", "成员2\n透明商品图"),
    ("02", "背景生成", "成员3\nSeedream 场景"),
    ("03", "无字底图", "成员4\n商品+背景素材"),
    ("04", "加文案", "成员4\n叠字成稿发布"),
]
for i, (num, title, body) in enumerate(steps):
    x = Inches(0.5 + i * 3.2)
    add_round(s, x, Inches(1.7), Inches(2.95), Inches(3.6), WHITE)
    add_text(s, x + Inches(0.2), Inches(1.95), Inches(2.5), Inches(0.4), num, size=26, bold=True, color=CYAN)
    add_text(s, x + Inches(0.2), Inches(2.55), Inches(2.5), Inches(0.4), title, size=18, bold=True, color=NAVY)
    add_text(s, x + Inches(0.2), Inches(3.2), Inches(2.5), Inches(1.4), body, size=14, color=MUTED)
    if i < 3:
        add_text(s, x + Inches(2.75), Inches(3.1), Inches(0.45), Inches(0.4), "→", size=20, bold=True, color=BLUE_MID)
add_bullets(s, Inches(0.5), Inches(5.6), Inches(12), Inches(1.0), [
    "海报合成拆成两步：先生成无字底图原始素材（可收藏），再基于底图加入营销文案成稿",
], size=13, color=MUTED)

# ===== 6 Features upper =====
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_chrome(s, "关键能力（上）", "选品 · 文案 · 抠图 · 背景")
left = [
    "智能选品：检索 ABO 货盘，一键跳转文案或海报",
    "文案生成：TikTok / Instagram / Amazon；中英日韩西",
    "多风格 + 批量多平台；历史可复用到海报",
]
right = [
    "商品抠图：上传 → 识别 → 自动抠图 → 展示下载",
    "背景生成：品类/风格/色调/场景/光照等 Prompt 控制",
    "与商品分析、文案、海报模块协同运行",
]
add_round(s, Inches(0.5), Inches(1.65), Inches(6.0), Inches(4.6), WHITE)
add_text(s, Inches(0.8), Inches(1.9), Inches(5.4), Inches(0.4), "选品与文案", size=18, bold=True, color=BLUE)
add_bullets(s, Inches(0.8), Inches(2.5), Inches(5.4), Inches(3.3), left, size=14)
add_round(s, Inches(6.8), Inches(1.65), Inches(6.0), Inches(4.6), WHITE)
add_text(s, Inches(7.1), Inches(1.9), Inches(5.4), Inches(0.4), "视觉前处理", size=18, bold=True, color=BLUE)
add_bullets(s, Inches(7.1), Inches(2.5), Inches(5.4), Inches(3.3), right, size=14)

# ===== 7 Features lower =====
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_chrome(s, "关键能力（下）", "海报 · 作品库 · 广场 · 客服看板")
feats = [
    ("海报模板与合成", "多文字层：字体/颜色/字号/坐标\n描边、阴影、发光等艺术字效果"),
    ("我的作品", "文案 · 成稿 · 底图素材 · 收藏\n底图可「用于加字」；成稿可发布"),
    ("作品广场", "浏览全站已发布海报\n搜索 / 预览 / 下载 / 收藏"),
    ("客服与看板", "会话推荐 + FAISS RAG\n运营数据、AI 建议、云端运维"),
]
for i, (h, body) in enumerate(feats):
    r, c = divmod(i, 2)
    x = Inches(0.5 + c * 6.35)
    y = Inches(1.6 + r * 2.4)
    add_round(s, x, y, Inches(6.1), Inches(2.2), WHITE)
    add_rect(s, x, y, Inches(0.1), Inches(2.2), CYAN if i % 2 == 0 else BLUE_MID)
    add_text(s, x + Inches(0.35), y + Inches(0.3), Inches(5.5), Inches(0.4), h, size=16, bold=True, color=BLUE)
    add_text(s, x + Inches(0.35), y + Inches(0.85), Inches(5.5), Inches(1.1), body, size=14, color=INK)

# ===== 8 Tech =====
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_chrome(s, "技术栈与 AI 能力", "工程可落地 · 外部模型可配置 · 无 Key 可 Mock")
tech = [
    ("前端", "Vue 3 · Vite · Element Plus · Pinia\nvue-i18n（5 语）· ECharts · Axios"),
    ("后端", "FastAPI · SQLAlchemy · SQLite · JWT\nPillow · Uvicorn"),
    ("大模型", "DeepSeek / OpenAI 兼容接口\n文案生成 · RAG 客服问答"),
    ("图像 AI", "火山方舟 Seedream 场景/精修\nrembg 抠图 · 可选 Ollama 视觉"),
]
for i, (h, body) in enumerate(tech):
    x = Inches(0.5 + (i % 2) * 6.35)
    y = Inches(1.65 + (i // 2) * 2.35)
    add_round(s, x, y, Inches(6.1), Inches(2.15), WHITE)
    add_rect(s, x, y, Inches(0.12), Inches(2.15), BLUE if i % 2 == 0 else CYAN)
    add_text(s, x + Inches(0.4), y + Inches(0.3), Inches(5.4), Inches(0.4), h, size=18, bold=True, color=NAVY)
    add_text(s, x + Inches(0.4), y + Inches(0.9), Inches(5.4), Inches(1.0), body, size=14, color=MUTED)

# ===== 9 Deploy =====
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_chrome(s, "演示与部署要点", "本地可跑通全链路；云端由成员5保障上线")
add_round(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(4.7), WHITE)
add_bullets(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.0), [
    "后端：pip install -r requirements.txt → python run.py → http://127.0.0.1:8000（Swagger /docs）",
    "前端：cd shared/frontend → npm install → npm run dev → http://localhost:5173",
    "环境：复制 .env.example → .env，配置 ARK_API_KEY、LLM/OPENAI、ABO 图片路径等",
    "账号：无内置演示账号，首次自行注册（用户名≥2、密码≥6）",
    "云端：构建发布、线上问题排查修复，保障系统稳定可用（成员5）",
], size=15)

# ===== 10 Summary =====
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_chrome(s, "总结与展望", "已打通主链路 · 差异化在「底图资产」与全站闭环")
add_round(s, Inches(0.5), Inches(1.65), Inches(7.8), Inches(4.7), WHITE)
add_bullets(s, Inches(0.85), Inches(2.0), Inches(7.2), Inches(4.0), [
    "已打通：选品 → 文案 / 海报四步 → 作品管理与广场 → RAG 客服与看板",
    "差异化：Seedream 场景 + 无字底图/叠字两段式素材资产",
    "工程亮点：模块化 monorepo、统一 JWT/API、Mock 便于无 Key 联调",
    "依赖提醒：Ark / DeepSeek Key、ABO 本地路径、rembg 首次拉模型",
], size=14)
add_round(s, Inches(8.55), Inches(1.65), Inches(4.25), Inches(4.7), NAVY)
add_text(s, Inches(8.9), Inches(2.1), Inches(3.6), Inches(0.45), "现场演示建议", size=16, bold=True, color=CYAN)
add_text(s, Inches(8.9), Inches(2.8), Inches(3.6), Inches(3.0),
         "1. 注册登录\n2. 智能选品\n3. 海报工作流四步\n4. 发布作品广场\n5. 打开运营看板",
         size=15, color=WHITE)

# ===== 11 Nav =====
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_chrome(s, "附录：产品导航地图", "侧栏入口与主路由对应关系")
nav = [
    ("首页", "/home"),
    ("智能选品中心", "/catalog"),
    ("文案生成", "/writing"),
    ("AI海报工作流", "/poster-workflow"),
    ("我的作品", "/my-works"),
    ("作品广场", "/gallery"),
    ("智能客服", "/chat"),
    ("运营看板", "/dashboard"),
]
for i, (name, path) in enumerate(nav):
    r, c = divmod(i, 4)
    x = Inches(0.5 + c * 3.2)
    y = Inches(1.7 + r * 2.3)
    add_round(s, x, y, Inches(3.0), Inches(2.0), WHITE)
    add_text(s, x + Inches(0.2), y + Inches(0.45), Inches(2.6), Inches(0.45), name, size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), y + Inches(1.1), Inches(2.6), Inches(0.4), path, size=12, color=BLUE_MID, align=PP_ALIGN.CENTER)

# ===== 12 Team overview =====
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_chrome(s, "团队协作分工总览", "五位成员分模块交付 · 统一接口联调")
overview = [
    ("成员1", "我的工作"),
    ("成员2", "商品识别与自动抠图\n标准化透明素材输出"),
    ("成员3", "背景生成与增强\nPrompt + Seedream"),
    ("成员4", "海报合成与素材管理\n模板 / 历史 / 收藏"),
    ("成员5", "客服·看板·选品\n云端部署与运维"),
]
for i, (h, body) in enumerate(overview):
    x = Inches(0.3 + i * 2.58)
    add_round(s, x, Inches(1.7), Inches(2.45), Inches(4.6), WHITE)
    add_rect(s, x, Inches(1.7), Inches(2.45), Inches(0.8), BLUE if i % 2 == 0 else NAVY)
    add_text(s, x + Inches(0.1), Inches(1.88), Inches(2.25), Inches(0.5), h, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), Inches(2.8), Inches(2.15), Inches(3.2), body, size=12, color=INK, align=PP_ALIGN.CENTER)

# ===== 13 M1 + M2 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_chrome(s, "分工详述：成员1 · 成员2", "")
member_card(
    s, Inches(0.45), Inches(1.55), Inches(6.1), Inches(4.9),
    "成员1 · 我的工作",
    "我的工作\n\n（此处可按实际补充本人负责模块、\n接口与演示亮点）",
)
member_card(
    s, Inches(6.8), Inches(1.55), Inches(6.1), Inches(4.9),
    "成员2 · 商品智能化处理",
    "负责商品图片智能化处理，打通：\n上传原图 → 识别商品 → 自动抠图 →\n结果展示 → 图片下载。\n\n为后续海报生成、营销文案等模块\n提供标准化的透明商品素材与识别数据。",
)

# ===== 14 M3 + M4 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_chrome(s, "分工详述：成员3 · 成员4", "")
member_card(
    s, Inches(0.45), Inches(1.55), Inches(6.1), Inches(4.9),
    "成员3 · 背景生成与增强",
    "负责 AI 海报工作流中「背景生成与增强」：\n承接成员2的商品理解与主体提取结果，\n基于类别、颜色、风格等结构化信息，\n完成商品信息 → Prompt 转换。\n\n控制场景、视觉风格、光照等参数，\n接入图像生成模型自动生成营销背景，\n并与商品分析、文案等功能协同。",
)
member_card(
    s, Inches(6.8), Inches(1.55), Inches(6.1), Inches(4.9),
    "成员4 · 海报合成与素材管理",
    "负责模块4：海报合成、模板管理、结果管理。\n将透明商品图 + 背景图 + 营销文案\n合成为可展示/下载/收藏的海报。\n\nFastAPI：上传、模板、合成、历史、\n下载统计、收藏；Pillow 多文字层与\n艺术字效果；模板/历史/收藏表持久化。",
)

# ===== 15 M5 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_chrome(s, "分工详述：成员5", "智能客服 · 运营看板 · 智能选品 · 云端部署")
add_round(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(4.9), WHITE)
add_rect(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.6), NAVY)
add_text(s, Inches(0.8), Inches(1.65), Inches(11.5), Inches(0.4),
         "成员5 · 客服 / 看板 / 选品 / 上线运维", size=16, bold=True, color=WHITE)
add_bullets(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(3.6), [
    "智能客服：完成客服会话与商品推荐展示",
    "运营看板：运营数据与 AI 建议呈现",
    "智能选品：选品浏览 / 活动推荐，并与文案、海报流程衔接",
    "云端部署：项目上线、构建发布、线上问题排查修复，保障系统稳定可用",
], size=15)

prs.save(OUT)
print(f"SAVED {OUT}")
